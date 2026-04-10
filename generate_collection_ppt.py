from __future__ import annotations

import json
import math
import sys
from collections import defaultdict
from pathlib import Path

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT / ".vendor"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUTPUT = ROOT / "Collection Dashboard.pptx"


COLORS = {
    "navy": RGBColor(36, 69, 112),
    "text": RGBColor(24, 24, 27),
    "muted": RGBColor(95, 102, 114),
    "line": RGBColor(220, 226, 234),
    "white": RGBColor(255, 255, 255),
    "blue_fill": RGBColor(230, 239, 250),
    "green_fill": RGBColor(226, 243, 236),
    "teal_fill": RGBColor(223, 241, 242),
    "blue": RGBColor(56, 114, 194),
    "green": RGBColor(47, 145, 98),
    "teal": RGBColor(42, 145, 153),
    "orange": RGBColor(225, 139, 16),
    "gray_bar": RGBColor(124, 124, 124),
    "dark_green": RGBColor(31, 122, 85),
    "soft_bg": RGBColor(249, 247, 244),
}


BUSINESS_COLUMNS = ["Step", "ITSC", "AA", "AB Service", "CD Service", "BJBP+CTC", "HK"]
BUSINESS_ROWS = [
    ["01 准备与客户主数据", "0.27", "0.01", "0.02", "0.01", "0.01", "0.04"],
    ["02 开票与其他支持材料交付", "10.50", "0.09", "0.22", "0.07", "0.02", "0.12"],
    ["03 对账与争议处理", "2.66", "0.21", "0.52", "0.63", "0.09", "0.21"],
    ["04 催收与承诺管理", "10.60", "0.57", "1.08", "1.10", "0.13", "0.72"],
    ["05 到账与核销", "2.22", "0.19", "0.13", "0.19", "0.09", "0.07"],
    ["06 票据/特殊收款", "0.00", "0.00", "0.00", "0.01", "0.00", "0.14"],
    ["07 升级与法律措施", "0.11", "0.01", "0.01", "0.01", "0.01", "-"],
    ["08 Write-off", "0.12", "0.00", "0.00", "0.01", "0.00", "0.02"],
    ["09 报表与持续改进", "0.72", "0.01", "0.02", "0.03", "0.00", "0.02"],
    ["10 其他", "0.24", "0.00", "0.03", "0.07", "0.02", "0.00"],
]
BUSINESS_TOTALS = ["Total FTE", "27.44", "1.09", "2.04", "2.12", "0.38", "1.33"]


def load_json(path: str):
    return json.loads((ROOT / path).read_text(encoding="utf-8-sig"))


def add_textbox(slide, x, y, w, h, text, font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT, font_name="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color or COLORS["text"]
    return box


def add_card(slide, x, y, w, h, fill, border, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1)
    return shape


def add_metric_card(slide, x, y, w, h, label, value, lines, fill, accent):
    add_card(slide, x, y, w, h, fill, COLORS["line"], radius=True)
    accent_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = accent
    accent_line.line.fill.background()
    add_textbox(slide, x + Inches(0.18), y + Inches(0.14), w - Inches(0.3), Inches(0.25), label.upper(), 10.5, True, COLORS["muted"])
    add_textbox(slide, x + Inches(0.18), y + Inches(0.42), w - Inches(0.3), Inches(0.45), value, 23, True, COLORS["navy"], "left", "Aptos Display")
    add_textbox(slide, x + Inches(0.18), y + Inches(0.92), w - Inches(0.3), Inches(0.55), "\n".join(lines), 11, False, COLORS["text"])


def add_rollout_card(slide, x, y, w, h, title, meta, actions, note, accent):
    add_card(slide, x, y, w, h, COLORS["white"], COLORS["line"], radius=True)
    accent_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.06), h)
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = accent
    accent_line.line.fill.background()
    add_textbox(slide, x + Inches(0.14), y + Inches(0.10), w - Inches(0.2), Inches(0.28), title, 15.5, True, accent, font_name="Aptos Display")
    add_textbox(slide, x + Inches(0.14), y + Inches(0.40), w - Inches(0.2), Inches(0.34), meta, 11.2, True, COLORS["text"])
    add_textbox(slide, x + Inches(0.14), y + Inches(0.78), Inches(0.72), Inches(0.28), "Actions:", 11.2, True, COLORS["muted"])
    add_textbox(slide, x + Inches(0.82), y + Inches(0.78), w - Inches(0.96), Inches(0.64), actions, 11.0, False, COLORS["text"])
    if note:
        add_textbox(slide, x + Inches(0.14), y + Inches(1.42), Inches(0.62), Inches(0.28), "Note:", 11.2, True, COLORS["muted"])
        add_textbox(slide, x + Inches(0.72), y + Inches(1.42), w - Inches(0.86), Inches(0.50), note, 10.8, False, COLORS["text"])


def to_num(value):
    if value in (None, "", "-"):
        return 0.0
    try:
        return float(value)
    except Exception:
        return 0.0


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(8.2), Inches(0.6), title, 26, True, COLORS["text"], font_name="Aptos Display")
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.95), Inches(8.4), Inches(0.35), subtitle, 12.5, False, COLORS["muted"])


def build_business_metrics(action_items):
    total_fte = sum(to_num(v) for v in BUSINESS_TOTALS[1:])
    target_fte = total_fte * 0.2
    assessment = sum(to_num(item.get("fte")) for item in action_items if item.get("status") == "评估中")
    pending_approval = sum(to_num(item.get("fte")) for item in action_items if item.get("status") == "待审批")
    approved = sum(to_num(item.get("fte")) for item in action_items if item.get("status") == "已审批")
    quote = sum(to_num(item.get("fte")) for item in action_items if item.get("status") == "待报价")
    fcst = assessment + pending_approval + approved + quote
    return {
        "total": total_fte,
        "target": target_fte,
        "fcst": fcst,
        "assessment": assessment,
        "approval": pending_approval,
        "approved": approved,
        "quote": quote,
    }


def normalize_step(value: str) -> str:
    return " ".join(str(value or "").split()).strip()


def build_process_summary(discussion, actions, threshold=0.1):
    merged = defaultdict(float)
    bands = {}
    for group in discussion.get("groups", []):
        for item in group.get("items", []):
            step = normalize_step(item.get("l2"))
            if not step:
                continue
            merged[(item.get("h") or group.get("l1") or "").strip(), step] += to_num(item.get("total"))
            bands[step] = (item.get("h") or group.get("l1") or "").strip()

    action_totals = defaultdict(float)
    linked_items = defaultdict(list)
    for item in actions.get("items", []):
        step = normalize_step(item.get("i") or item.get("step"))
        if not step:
            continue
        action_totals[step] += to_num(item.get("fte"))
        linked_items[step].append(item)

    ordered = []
    for (band, step), total in merged.items():
        if total + 1e-9 < threshold:
            continue
        ordered.append({
            "band": band or "Other",
            "step": step,
            "total": total,
            "action_total": action_totals.get(step, 0.0),
            "linked_items": linked_items.get(step, []),
        })
    return ordered


def add_business_table(slide, x, y, w, h):
    rows = len(BUSINESS_ROWS) + 2
    cols = len(BUSINESS_COLUMNS)
    table = slide.shapes.add_table(rows, cols, x, y, w, h).table
    widths = [2.4, 0.92, 0.92, 0.92, 0.92, 0.92, 0.92]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)

    for idx, header in enumerate(BUSINESS_COLUMNS):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(242, 243, 245)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if idx else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(10)
                run.font.color.rgb = COLORS["muted"]

    for r_idx, row in enumerate(BUSINESS_ROWS, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(10.5)
                    run.font.color.rgb = COLORS["text"] if c_idx != 1 else RGBColor(209, 103, 47)
                    run.font.bold = c_idx == 1

    last = rows - 1
    for c_idx, val in enumerate(BUSINESS_TOTALS):
        cell = table.cell(last, c_idx)
        cell.text = str(val)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(248, 248, 248)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(10.5)
                run.font.bold = True
                run.font.color.rgb = COLORS["text"]
    return table


def add_kpi_panel(slide, x, y, w, h, title, value, note=None):
    add_card(slide, x, y, w, h, COLORS["white"], COLORS["line"], radius=True)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.16), w - Inches(0.3), Inches(0.22), title, 12, False, COLORS["muted"])
    add_textbox(slide, x + Inches(0.18), y + Inches(0.42), w - Inches(0.3), Inches(0.42), value, 24, True, COLORS["text"], font_name="Aptos Display")
    if note:
        add_textbox(slide, x + Inches(0.18), y + Inches(0.88), w - Inches(0.3), h - Inches(0.98), note, 10.8, False, COLORS["muted"])


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    discussion = load_json("discussion-summary.json")
    actions = load_json("action-summary.json")
    action_items = actions.get("items", [])
    process_items = build_process_summary(discussion, actions, 0.1)
    business_metrics = build_business_metrics(action_items)

    status_rank = {"待评估": 0, "评估中": 1, "待报价": 2, "待审批": 3, "已审批": 4, "取消": 5}
    sorted_actions = sorted(
        action_items,
        key=lambda item: (status_rank.get(item.get("status"), 99), -to_num(item.get("fte"))),
    )

    # Slide 1: executive summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "2026–2027 Accounting Cost Reduction Plan", "This plan aims to reduce Accounting operating cost through a combination of GBS migration and AI/Automation-led process optimization across AP/AR, GL, Collection and Invoicing.")
    add_metric_card(slide, Inches(0.45), Inches(1.55), Inches(4.0), Inches(1.38), "Savings Potential", "$0.54M - 1.23M", ["Low case without GBS: $544.2K", "High case with GBS: $1.12M"], COLORS["blue_fill"], COLORS["blue"])
    add_metric_card(slide, Inches(4.65), Inches(1.55), Inches(4.0), Inches(1.38), "Investment & Return", "$150K / 0.8 - 2.6 yrs", ["AI & Automation investment: $150K", "ROI range: 0.8 to 2.6 years"], COLORS["green_fill"], COLORS["green"])
    add_metric_card(slide, Inches(8.85), Inches(1.55), Inches(4.0), Inches(1.38), "Impact Application", "5 + 3", ["Existing: PCS / VAT / AR Portal / MDM / Credit", "New: Enterprise WeChat / AI / RPA"], COLORS["teal_fill"], COLORS["teal"])
    add_rollout_card(slide, Inches(0.42), Inches(3.25), Inches(6.1), Inches(1.58), "1. AP / AR", "Kickoff: Apr 2026   |   Existed Cost: $1,715K   |   Saving Target: $520K", "Total HC 34, convert the AP/AR team to the GBS model; keep 8 HC in Shanghai and move 23 HC to GBS.", "$20K / HC or 50% of China Cost, 9 months transition cycle", COLORS["blue"])
    add_rollout_card(slide, Inches(6.68), Inches(3.25), Inches(6.1), Inches(1.58), "2. Collection", "Kickoff: April 2026   |   Existed Cost: $1,570K   |   Saving Target: $314K", "Total HC 36, deliver around 20% savings through AI/Automation-enabled productivity improvement.", "Detailed scope to be identified.", COLORS["teal"])
    add_rollout_card(slide, Inches(0.42), Inches(4.95), Inches(6.1), Inches(1.45), "3. GL", "Kickoff: July 2026   |   Existed Cost: $1,923K   |   Saving Target: $219K", "Total HC 28, convert part of the GL team (14 FTE) to the GBS model after process mapping and operating-model validation.", "100% band 5 and 30% band 6 move to GBS", COLORS["green"])
    add_rollout_card(slide, Inches(6.68), Inches(4.95), Inches(6.1), Inches(1.45), "4. Invoicing", "Kickoff: April 2026   |   Existed Cost: $398K   |   Saving Target: $70K", "Total HC 9, leverage AI/Automation to improve invoicing efficiency by about 20%.", "Keeping local execution for compliance", COLORS["orange"])

    # Slide 2: collection overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Collection Optimization Leverage AI/Automation", "The program is focused on collection workflow optimization, cross-team automation enablement, and rollout readiness. Overall momentum is stable, with the main constraints concentrated around integration bandwidth and stakeholder sign-off.")
    add_card(slide, Inches(0.55), Inches(1.65), Inches(6.1), Inches(4.95), COLORS["white"], COLORS["line"], radius=True)
    add_textbox(slide, Inches(0.82), Inches(2.05), Inches(2.8), Inches(1.65), "20%", 48, True, COLORS["text"], font_name="Aptos Display")
    add_textbox(slide, Inches(0.9), Inches(3.15), Inches(2.6), Inches(0.3), "Overall Progress", 12, False, COLORS["muted"])
    for idx, (label, value) in enumerate([("Project Timeline", "2026.02.01 - 2026.06.30"), ("Current Phase", "Assessment"), ("Owner", "Heather Sun")]):
        card_x = Inches(3.5 + idx * 1.02)
        add_card(slide, card_x, Inches(2.15), Inches(0.95), Inches(1.7), COLORS["soft_bg"], COLORS["line"], radius=True)
        add_textbox(slide, card_x + Inches(0.1), Inches(2.33), Inches(0.78), Inches(0.26), label, 9.5, False, COLORS["muted"])
        add_textbox(slide, card_x + Inches(0.1), Inches(2.72), Inches(0.78), Inches(0.7), value, 12.5, True, COLORS["text"])
    add_textbox(slide, Inches(7.05), Inches(1.78), Inches(5.8), Inches(0.4), "Program Overview", 20, True, COLORS["navy"], font_name="Aptos Display")
    bullets = [
        "Assessment stage complete enough to frame the cost-reduction path.",
        "Collection is positioned as the main AI / Automation productivity lever.",
        "Execution focus is now on scope confirmation, prioritization, and delivery planning.",
    ]
    for idx, bullet in enumerate(bullets):
        add_textbox(slide, Inches(7.1), Inches(2.35 + idx * 0.64), Inches(5.2), Inches(0.5), f"• {bullet}", 14, False, COLORS["text"])

    # Slide 3: business current status
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Business Current Status", "Current FTE allocation by business step and legal entity, based on the latest workload snapshot.")
    add_business_table(slide, Inches(0.45), Inches(1.4), Inches(8.35), Inches(5.35))
    add_kpi_panel(slide, Inches(9.05), Inches(1.4), Inches(3.7), Inches(1.0), "Total FTE", f"{business_metrics['total']:.3f}")
    add_kpi_panel(slide, Inches(9.05), Inches(2.55), Inches(3.7), Inches(1.15), "Target FTE", f"{business_metrics['target']:.3f}", "20% of Total FTE")
    add_card(slide, Inches(9.05), Inches(3.9), Inches(3.7), Inches(2.85), COLORS["white"], COLORS["line"], radius=True)
    add_textbox(slide, Inches(9.25), Inches(4.08), Inches(1.5), Inches(0.25), "FCST FTE", 12, False, COLORS["muted"])
    add_textbox(slide, Inches(9.25), Inches(4.34), Inches(1.8), Inches(0.4), f"{business_metrics['fcst']:.3f}", 24, True, COLORS["text"], font_name="Aptos Display")
    labels = [("评估中", business_metrics["assessment"]), ("待审批", business_metrics["approval"]), ("已审批", business_metrics["approved"]), ("待报价", business_metrics["quote"])]
    y = 5.1
    for label, value in labels:
        add_textbox(slide, Inches(9.28), Inches(y), Inches(1.5), Inches(0.22), label, 11.5, False, COLORS["text"])
        add_textbox(slide, Inches(11.85), Inches(y), Inches(0.7), Inches(0.22), f"{value:.3f}", 11.5, True, COLORS["muted"], PP_ALIGN.RIGHT)
        y += 0.36

    # Slide 4: process footprint
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Process Footprint View")
    add_textbox(slide, Inches(9.4), Inches(0.45), Inches(1.25), Inches(0.24), "FTE Threshold", 11.5, False, COLORS["muted"])
    add_card(slide, Inches(10.55), Inches(0.33), Inches(1.1), Inches(0.42), COLORS["white"], COLORS["line"], radius=True)
    add_textbox(slide, Inches(10.72), Inches(0.42), Inches(0.72), Inches(0.18), "0.10", 12, False, COLORS["text"])
    add_textbox(slide, Inches(9.4), Inches(0.78), Inches(2.2), Inches(0.22), "Flows below this value are hidden.", 9.5, False, COLORS["muted"])
    add_textbox(slide, Inches(0.72), Inches(1.0), Inches(1.8), Inches(0.22), "Gray: No optimization potential", 10.5, False, COLORS["muted"])
    add_textbox(slide, Inches(2.95), Inches(1.0), Inches(2.0), Inches(0.22), "Green: Click for detailed actions", 10.5, False, COLORS["muted"])
    baseline_y = 3.2
    start_x = 0.82
    usable_w = 11.8
    count = max(len(process_items), 1)
    step_w = usable_w / count
    max_total = max((item["total"] for item in process_items), default=0.1)
    prev_band = None
    band_start = 0
    for idx, item in enumerate(process_items + [{"band": None}]):
        if idx == len(process_items) or item["band"] != prev_band:
            if prev_band is not None:
                x = Inches(start_x + band_start * step_w)
                w = Inches((idx - band_start) * step_w)
                band = add_card(slide, x, Inches(2.45), w, Inches(0.34), COLORS["teal_fill"] if band_start % 2 == 0 else COLORS["blue_fill"], COLORS["line"])
                add_textbox(slide, x, Inches(2.49), w, Inches(0.18), prev_band, 11.5, True, COLORS["text"], PP_ALIGN.CENTER)
            band_start = idx
            prev_band = item["band"] if idx < len(process_items) else None
    for idx, item in enumerate(process_items):
        cx = start_x + idx * step_w + step_w * 0.52
        label = slide.shapes.add_textbox(Inches(cx), Inches(2.38), Inches(0.15), Inches(1.1))
        tf = label.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = item["step"]
        run.font.size = Pt(8.5)
        run.font.color.rgb = RGBColor(60, 60, 60)
        run.font.name = "Noto Sans SC"
        label.rotation = 270
        bar_h = max((item["total"] / max_total) * 2.15, 0.12)
        bar_x = Inches(start_x + idx * step_w + step_w * 0.42)
        bar_y = Inches(baseline_y)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, bar_x, bar_y, Inches(min(step_w * 0.32, 0.22)), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS["dark_green"] if item["action_total"] > 0 else COLORS["gray_bar"]
        bar.line.fill.background()
        add_textbox(slide, Inches(start_x + idx * step_w + step_w * 0.22), Inches(baseline_y + bar_h + 0.02), Inches(step_w * 0.6), Inches(0.2), f"{item['total']:.1f}", 9, False, COLORS["muted"], PP_ALIGN.CENTER)

    # Slide 5: improvement points
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Identified Improvement Points", "Current change opportunities summarized from the workbook action sheet.")
    summary = defaultdict(float)
    for item in action_items:
        summary[item.get("status") or "Unknown"] += to_num(item.get("fte"))
    status_cards = [("待评估", summary["待评估"]), ("评估中", summary["评估中"]), ("待报价", summary["待报价"]), ("待审批", summary["待审批"])]
    for idx, (label, value) in enumerate(status_cards):
        x = Inches(0.48 + idx * 3.12)
        add_kpi_panel(slide, x, Inches(1.25), Inches(2.84), Inches(0.95), label, f"{value:.3f}")
    card_specs = [
        (Inches(0.48), Inches(2.45), Inches(4.1), Inches(2.0)),
        (Inches(4.63), Inches(2.45), Inches(4.1), Inches(2.0)),
        (Inches(8.78), Inches(2.45), Inches(4.1), Inches(2.0)),
        (Inches(0.48), Inches(4.65), Inches(4.1), Inches(2.0)),
        (Inches(4.63), Inches(4.65), Inches(4.1), Inches(2.0)),
        (Inches(8.78), Inches(4.65), Inches(4.1), Inches(2.0)),
    ]
    for item, (x, y, w, h) in zip(sorted_actions[:6], card_specs):
        tone = COLORS["dark_green"] if to_num(item.get("fte")) > 0 else COLORS["gray_bar"]
        add_card(slide, x, y, w, h, COLORS["white"], COLORS["line"], radius=True)
        top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.05))
        top.fill.solid()
        top.fill.fore_color.rgb = tone
        top.line.fill.background()
        add_textbox(slide, x + Inches(0.18), y + Inches(0.16), w - Inches(0.36), Inches(0.45), item.get("changePoint") or "-", 16, True, COLORS["text"], font_name="Aptos Display")
        add_textbox(slide, x + Inches(0.18), y + Inches(0.72), Inches(1.12), Inches(0.24), f"Type: {item.get('type') or '-'}", 10.8, True, COLORS["text"])
        add_textbox(slide, x + Inches(1.7), y + Inches(0.72), Inches(0.95), Inches(0.24), f"FTE: {to_num(item.get('fte')):.3f}", 10.8, True, COLORS["text"])
        cost = item.get("cost")
        cost_text = "-" if cost in (None, "", "-") else f"￥{int(float(cost)):,}"
        add_textbox(slide, x + Inches(2.8), y + Inches(0.72), Inches(1.05), Inches(0.24), f"Cost: {cost_text}", 10.8, True, COLORS["text"])
        add_textbox(slide, x + Inches(0.18), y + Inches(1.05), w - Inches(0.36), Inches(0.76), (item.get("remark") or "-").split("\n")[0], 10.5, False, COLORS["muted"])

    prs.save(OUTPUT)
    print(f"Saved PPT to {OUTPUT}")


if __name__ == "__main__":
    create_presentation()
